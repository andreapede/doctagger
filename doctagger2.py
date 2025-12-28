import os
import argparse
import json
import datetime
import platform
import ollama

from pypdf import PdfReader
from docx import Document
from pptx import Presentation


class DocumentAnalyzer:
    def __init__(self, filepath, model="gemma3", dump_chunks=False, debug_reduce=False):
        self.filepath = filepath
        self.filename = os.path.basename(filepath)
        self.extension = os.path.splitext(filepath)[1].lower()
        self.model = model
        self.content = ""
        self.metadata = {}
        self.dump_chunks = dump_chunks
        self.debug_reduce = debug_reduce

    def _ollama_json_chat(self, prompt, options=None, retries=2):
        """Call Ollama expecting JSON and parse safely, with low temperature and retries."""
        for attempt in range(retries):
            try:
                response = ollama.chat(
                    model=self.model,
                    messages=[{"role": "user", "content": prompt}],
                    format="json",
                    options=options or {"temperature": 0},
                )

                content = response.get("message", {}).get("content", "")
                data = json.loads(content) if content else {}
                if isinstance(data, dict) and data:
                    return data
            except Exception:
                pass
        return {}

    def _extract_json_object(self, text):
        """Attempt to extract a JSON object from free-form text by bracket matching."""
        if not text:
            return {}
        start = text.find("{")
        while start != -1:
            depth = 0
            for i in range(start, len(text)):
                ch = text[i]
                if ch == "{":
                    depth += 1
                elif ch == "}":
                    depth -= 1
                    if depth == 0:
                        candidate = text[start : i + 1]
                        try:
                            obj = json.loads(candidate)
                            if isinstance(obj, dict):
                                return obj
                        except Exception:
                            break
            start = text.find("{", start + 1)
        return {}

    def process(self):
        print(f"[*] Processing file: {self.filename}...")
        self._extract_system_metadata()

        try:
            self._extract_text()
        except Exception as e:
            print(f"[!] Text extraction error: {e}")
            return

        if not self.content:
            print("[!] The file appears empty or unreadable.")
            return

        print(f"[*] AI analysis in progress with local model '{self.model}'... (this may take time)")
        self._analyze_content_ai()

        self._save_json()
        print(f"[V] Done. Metadata saved in {self.filename}.json")

    def _extract_system_metadata(self):
        stats = os.stat(self.filepath)
        if platform.system() == "Windows":
            creation_time = stats.st_ctime
        else:
            try:
                creation_time = stats.st_birthtime
            except AttributeError:
                creation_time = stats.st_ctime

        self.metadata["file_info"] = {
            "filename": self.filename,
            "extension": self.extension,
            "size_bytes": stats.st_size,
            "created_at": datetime.datetime.fromtimestamp(creation_time).isoformat(),
            "modified_at": datetime.datetime.fromtimestamp(stats.st_mtime).isoformat(),
        }

    def _extract_text(self):
        if self.extension == ".txt":
            with open(self.filepath, "r", encoding="utf-8") as f:
                self.content = f.read()
        elif self.extension == ".pdf":
            reader = PdfReader(self.filepath)
            text_parts = [page.extract_text() or "" for page in reader.pages]
            self.content = "\n".join(text_parts)
        elif self.extension == ".docx":
            doc = Document(self.filepath)
            self.content = "\n".join([para.text for para in doc.paragraphs])
        elif self.extension == ".pptx":
            prs = Presentation(self.filepath)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
            self.content = "\n".join(text_parts)

        self.content = self.content.strip()

    def _analyze_content_ai(self):
        print(f"[*] Avvio analisi profonda (Map-Reduce) su {len(self.content)} caratteri...")

        # 1. CHUNKING
        chunks = self._smart_chunking(self.content, chunk_size=6000, overlap=500)
        print(f"[*] Documento diviso in {len(chunks)} parti. Generazione riassunti intermedi...")

        partial_summaries = []
        partials = []

        # 2. MAP (Sintesi Parziale)
        for i, chunk in enumerate(chunks):
            map_prompt = f"""
            Analizza questo frammento di testo (Parte {i + 1} di {len(chunks)}).
            Estrai i punti chiave tecnici, le metodologie usate e i risultati descritti.
            Sii conciso. Non usare preamboli.

            TESTO:
            {chunk}
            """

            try:
                response = ollama.chat(
                    model=self.model, messages=[{"role": "user", "content": map_prompt}]
                )
                summary = response.get("message", {}).get("content", "")
                if summary:
                    partial_summaries.append(summary)
                print(f"    - Chunk {i + 1}/{len(chunks)} analizzato.")
                partials.append({
                    "index": i,
                    "chunk_chars": len(chunk),
                    "summary": summary,
                    "error": None,
                    "model": self.model,
                    "timestamp": datetime.datetime.now().isoformat(),
                })
            except Exception as e:
                print(f"    [!] Errore chunk {i}: {e}")
                partials.append({
                    "index": i,
                    "chunk_chars": len(chunk),
                    "summary": "",
                    "error": str(e),
                    "model": self.model,
                    "timestamp": datetime.datetime.now().isoformat(),
                })

        if self.dump_chunks:
            self._save_chunks_json(partials)

        # 3. REDUCE (Sintesi Finale) con riduzione gerarchica per evitare superamento contesto
        print("[*] Generazione Abstract Finale e Tag...")

        def hierarchical_reduce(summaries):
            if not summaries:
                return {}

            group_size = 6
            groups = [summaries[i : i + group_size] for i in range(0, len(summaries), group_size)]

            condensed = []
            for gi, grp in enumerate(groups):
                grp_text = "\n---\n".join(grp)
                grp_prompt = f"""
                Riassumi questi riassunti intermedi (Gruppo {gi + 1}) in un paragrafo coeso in Inglese.
                OUTPUT JSON:
                {{"abstract": "..."}}

                INPUT:
                {grp_text}
                """
                data = self._ollama_json_chat(grp_prompt)
                if data.get("abstract"):
                    condensed.append(data["abstract"])
                elif self.debug_reduce:
                    # Try to capture raw output for debugging
                    try:
                        raw = ollama.chat(
                            model=self.model,
                            messages=[{"role": "user", "content": grp_prompt}],
                            options={"temperature": 0},
                        )
                        raw_text = raw.get("message", {}).get("content", "")
                        with open(f"{self.filename}.reduce_group_{gi}.raw.txt", "w", encoding="utf-8") as rf:
                            rf.write(raw_text)
                    except Exception:
                        pass

            final_input = "\n---\n".join(condensed) if condensed else "\n---\n".join(summaries)

            reduce_prompt = f"""
            Sei un analista tecnico esperto. Qui di seguito trovi una serie di riassunti estratti da un documento tecnico.

            Il tuo compito è generare un output JSON strutturato basandoti sull'UNIONE di queste informazioni.

            REGOLE OUTPUT JSON:
            1. 'abstract': Un testo coeso e scorrevole (NON un elenco puntato) di circa 300-400 parole che descriva: Scopo, Metodologia (dettagli tecnici), Risultati e Conclusioni.
            2. 'tags_it': 10 tag tecnici specifici in Italiano.
            3. 'tags_en': 10 tag tecnici specifici in Inglese.
            4. 'technical_specs': Una lista di strumenti, software (es. LabVIEW), o algoritmi citati.

            RIASSUNTI:
            {final_input}
            """

            # First: ask for JSON directly
            direct = self._ollama_json_chat(reduce_prompt, options={"temperature": 0})
            if direct:
                return direct

            # Second: get raw text and try to extract JSON
            try:
                raw = ollama.chat(
                    model=self.model,
                    messages=[{"role": "user", "content": reduce_prompt}],
                    options={"temperature": 0},
                )
                raw_text = raw.get("message", {}).get("content", "")
            except Exception:
                raw_text = ""

            if self.debug_reduce and raw_text:
                with open(f"{self.filename}.reduce_raw.txt", "w", encoding="utf-8") as f:
                    f.write(raw_text)

            extracted = self._extract_json_object(raw_text)
            if extracted:
                return extracted

            return {}

        def fallback_reduce(summaries):
            """Fallback: ask for fields one by one and merge."""
            base_text = "\n---\n".join(summaries)

            ap = f"""
            Write ONLY JSON with key 'abstract' (English ~300-400 words) summarizing scope, methodology (technical details), results, conclusions.
            INPUT:
            {base_text}
            """
            tp_it = f"""
            Return ONLY JSON with key 'tags_it' as an array of 10 specific technical tags in Italian.
            INPUT:
            {base_text}
            """
            tp_en = f"""
            Return ONLY JSON with key 'tags_en' as an array of 10 specific technical tags in English.
            INPUT:
            {base_text}
            """
            sp = f"""
            Return ONLY JSON with key 'technical_specs' as an array listing tools/software/algorithms cited.
            INPUT:
            {base_text}
            """

            result = {}
            for name, prompt in [("abstract", ap), ("tags_it", tp_it), ("tags_en", tp_en), ("technical_specs", sp)]:
                data = self._ollama_json_chat(prompt, options={"temperature": 0})
                if not data:
                    try:
                        raw = ollama.chat(model=self.model, messages=[{"role": "user", "content": prompt}], options={"temperature": 0})
                        raw_text = raw.get("message", {}).get("content", "")
                        data = self._extract_json_object(raw_text)
                    except Exception:
                        data = {}
                if name in data:
                    result[name] = data[name]

            return result

        try:
            data = hierarchical_reduce(partial_summaries)
            if not data:
                data = fallback_reduce(partial_summaries)

            self.metadata["analysis"] = {
                "abstract": data.get("abstract", "N/A"),
                "tags": {
                    "it": data.get("tags_it", []),
                    "en": data.get("tags_en", []),
                },
                "technical_specs": data.get("technical_specs", []),
                "method": "map_reduce",
                "chunks_processed": len(chunks),
                "word_count": len(self.content.split()),
                "model_used": self.model,
                "processed_date": datetime.datetime.now().isoformat(),
                "partials_dump": self.dump_chunks,
                "reduce_debug": self.debug_reduce,
                "partial_chunks": partials,
            }

            # Se l'output è vuoto, registra un avviso
            if self.metadata["analysis"]["abstract"] == "N/A" and not (
                self.metadata["analysis"]["tags"]["it"] or self.metadata["analysis"]["tags"]["en"]
            ):
                print("[!] Output JSON vuoto dal Reduce. Considera pull/aggiornare il modello o riprovare.")

        except Exception as e:
            print(f"[!] Errore nel Reduce finale: {e}")
            self.metadata["analysis"] = {
                "partial_summaries": partial_summaries,
                "error": str(e),
            }

    def _save_chunks_json(self, partials):
        output_filename = f"{self.filename}.chunks.json"
        with open(output_filename, "w", encoding="utf-8") as f:
            json.dump({
                "file": self.filename,
                "model": self.model,
                "chunks": partials,
            }, f, indent=4, ensure_ascii=False)
        print(f"[*] Dump dei chunk salvato in {output_filename}")

    def _smart_chunking(self, text, chunk_size=6000, overlap=500):
        """Divide il testo in chunk con sovrapposizione."""
        if len(text) <= chunk_size:
            return [text]

        chunks = []
        start = 0
        text_len = len(text)

        while start < text_len:
            end = start + chunk_size

            if end >= text_len:
                chunks.append(text[start:])
                break

            snippet = text[end - 100 : end]
            last_space = snippet.rfind(" ")

            if last_space != -1:
                end = (end - 100) + last_space

            chunks.append(text[start:end])

            start = end - overlap

        return chunks

    def _save_json(self):
        output_filename = f"{self.filename}.json"
        with open(output_filename, "w", encoding="utf-8") as f:
            json.dump(self.metadata, f, indent=4, ensure_ascii=False)


if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        prog="doctagger2.py",
        description="DocTagger – AI-powered document analysis using local Ollama models.",
        epilog="""
Examples:
  # Basic usage with default model (gemma3)
  python doctagger2.py document.pdf
  
  # Use a different model
  python doctagger2.py document.pdf --model mistral
  
  # Save per-chunk summaries
  python doctagger2.py document.pdf --dump-chunks
  
  # Enable reduce debugging (for troubleshooting empty output)
  python doctagger2.py document.pdf --debug-reduce
  
  # Combine multiple flags
  python doctagger2.py document.pdf --model gemma3 --dump-chunks --debug-reduce

Supported formats: PDF, DOCX, PPTX, TXT

Output files:
  - <filename>.json – Main analysis with abstract, tags (IT/EN), technical specs, and partial chunks.
  - <filename>.chunks.json – Per-chunk summaries (with --dump-chunks).
  - <filename>.reduce_raw.txt – Raw reduce output (with --debug-reduce if JSON parsing failed).
  - <filename>.reduce_group_<N>.raw.txt – Intermediate group outputs (with --debug-reduce).

For more information, see README.md or run:
  python doctagger2.py --help
        """,
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument(
        "file",
        help="Path to the document to analyze (PDF, DOCX, PPTX, or TXT)."
    )
    parser.add_argument(
        "--model",
        default="gemma3",
        help="Ollama model to use for analysis. Default: gemma3. Examples: mistral, llama2, llama3.2."
    )
    parser.add_argument(
        "--dump-chunks",
        action="store_true",
        help="Save per-chunk analysis (index, char count, summary, error, timestamp) to <filename>.chunks.json."
    )
    parser.add_argument(
        "--debug-reduce",
        action="store_true",
        help="Dump raw reduce outputs and enable JSON extraction fallback. Writes <filename>.reduce_raw.txt and optionally <filename>.reduce_group_*.raw.txt."
    )
    args = parser.parse_args()

    if os.path.exists(args.file):
        analyzer = DocumentAnalyzer(
            args.file,
            model=args.model,
            dump_chunks=args.dump_chunks,
            debug_reduce=args.debug_reduce,
        )
        analyzer.process()
    else:
        print("File not found.")