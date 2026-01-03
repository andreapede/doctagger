import os
import argparse
import json
import datetime
import platform
import ollama

from collections import Counter, defaultdict

from pypdf import PdfReader
from docx import Document
from pptx import Presentation


class DocumentAnalyzer:
    def __init__(self, filepath, model="gemma3", dump_chunks=False, debug_reduce=False):
        self.filepath = filepath
        self.filename = os.path.basename(filepath)
        self.extension = os.path.splitext(filepath)[1].lower()
        self.model = model
        self.content = ""
        self.metadata = {}
        self.dump_chunks = dump_chunks
        self.debug_reduce = debug_reduce

    def _ollama_json_chat(self, prompt, options=None, retries=2):
        """Call Ollama expecting JSON and parse safely, with low temperature and retries."""
        for attempt in range(retries):
            try:
                response = ollama.chat(
                    model=self.model,
                    messages=[{"role": "user", "content": prompt}],
                    format="json",
                    options=options or {"temperature": 0},
                )

                content = response.get("message", {}).get("content", "")
                try:
                    data = json.loads(content) if content else {}
                except json.JSONDecodeError:
                    data = self._extract_json_object(content)

                if isinstance(data, dict) and data:
                    return data
            except Exception as e:
                if self.debug_reduce:
                    print(f"    [DEBUG] _ollama_json_chat attempt {attempt+1} error: {e}")
                pass
        return {}

    def _normalize_tags_field(self, v):
        """Normalize a tags/spec field to list[str]. Accepts list, str, dict-indexed."""
        if v is None:
            return []
        if isinstance(v, list):
            return [str(x).strip() for x in v if str(x).strip()]
        if isinstance(v, dict):
            # dict indicizzato {"0":"a","1":"b"} o simili
            if all(str(k).isdigit() for k in v.keys()):
                return [str(v[k]).strip() for k in sorted(v.keys(), key=lambda x: int(x)) if str(v[k]).strip()]
            return []
        if isinstance(v, str):
            parts = []
            for chunk in v.replace(";", "\n").replace(",", "\n").splitlines():
                s = chunk.strip()
                if s:
                    parts.append(s)
            return parts
        s = str(v).strip()
        return [s] if s else []

    def _rank_union(self, list_of_lists, top_n=10):
        """Union + ranking by frequency (case-insensitive), preserving most common display form."""
        counts = Counter()
        display = defaultdict(Counter)

        for tags in list_of_lists:
            for t in self._normalize_tags_field(tags):
                key = t.lower().strip()
                if not key:
                    continue
                counts[key] += 1
                display[key][t] += 1

        ranked_keys = sorted(counts.keys(), key=lambda k: (-counts[k], k))

        out = []
        for k in ranked_keys:
            best = display[k].most_common(1)[0][0]
            out.append(best)
            if len(out) >= top_n:
                break
        return out

    def _deduplicate_tags_llm(self, tags, language="it"):
        """Use LLM to remove semantic duplicates from a list of tags."""
        if not tags:
            return []

        prompt = f"""
You are a terminologist.
Clean up this list of {language.upper()} tags by removing semantic duplicates (e.g. singular/plural, synonyms).
Keep the most precise or common form.
Return ONLY a JSON object with a single key "tags" containing the cleaned list of strings.

INPUT TAGS:
{json.dumps(tags, ensure_ascii=False)}

OUTPUT JSON:
{{ "tags": ["tag1", "tag2", ...] }}
"""
        data = self._ollama_json_chat(prompt)
        cleaned = self._normalize_tags_field(data.get("tags", []))

        if not cleaned:
            return tags

        return cleaned

    def _extract_json_object(self, text):
        """Attempt to extract a JSON object from free-form text by bracket matching."""
        if not text:
            return {}
        start = text.find("{")
        while start != -1:
            depth = 0
            for i in range(start, len(text)):
                ch = text[i]
                if ch == "{":
                    depth += 1
                elif ch == "}":
                    depth -= 1
                    if depth == 0:
                        candidate = text[start : i + 1]
                        try:
                            obj = json.loads(candidate)
                            if isinstance(obj, dict):
                                return obj
                        except Exception:
                            break
            start = text.find("{", start + 1)
        return {}

    def process(self):
        print(f"[*] Processing file: {self.filename}...")
        self._extract_system_metadata()

        try:
            self._extract_text()
        except Exception as e:
            print(f"[!] Text extraction error: {e}")
            return

        if not self.content:
            print("[!] The file appears empty or unreadable.")
            return

        print(f"[*] AI analysis in progress with local model '{self.model}'... (this may take time)")
        self._analyze_content_ai()

        self._save_json()
        print(f"[V] Done. Metadata saved in {self.filename}.json")

    def _extract_system_metadata(self):
        stats = os.stat(self.filepath)
        if platform.system() == "Windows":
            creation_time = stats.st_ctime
        else:
            try:
                creation_time = stats.st_birthtime
            except AttributeError:
                creation_time = stats.st_ctime

        self.metadata["file_info"] = {
            "filename": self.filename,
            "extension": self.extension,
            "size_bytes": stats.st_size,
            "created_at": datetime.datetime.fromtimestamp(creation_time).isoformat(),
            "modified_at": datetime.datetime.fromtimestamp(stats.st_mtime).isoformat(),
        }

    def _extract_text(self):
        if self.extension == ".txt":
            with open(self.filepath, "r", encoding="utf-8") as f:
                self.content = f.read()
        elif self.extension == ".pdf":
            reader = PdfReader(self.filepath)
            text_parts = [page.extract_text() or "" for page in reader.pages]
            self.content = "\n".join(text_parts)
        elif self.extension == ".docx":
            doc = Document(self.filepath)
            self.content = "\n".join([para.text for para in doc.paragraphs])
        elif self.extension == ".pptx":
            prs = Presentation(self.filepath)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_parts.append(shape.text)
            self.content = "\n".join(text_parts)

        self.content = self.content.strip()

    def _analyze_content_ai(self):
        print(f"[*] Avvio analisi profonda (Map-Reduce) su {len(self.content)} caratteri...")

        # 1. CHUNKING
        chunks = self._smart_chunking(self.content, chunk_size=6000, overlap=500)
        print(f"[*] Documento diviso in {len(chunks)} parti. Generazione riassunti intermedi...")

        partial_summaries = []
        partials = []

        
        # 2. MAP (Analisi Parziale) - JSON per chunk (summary + tag IT/EN + technical specs)
        for i, chunk in enumerate(chunks):
            map_prompt = f"""
Sei un analista tecnico.
Analizza questo frammento (Parte {i + 1} di {len(chunks)}).

RESTITUISCI SOLO JSON (niente testo fuori dal JSON) con questo schema:
{{
  "summary": "max 8-10 righe, inglese, punti chiave tecnici e risultati",
  "tags_it": ["6-10 tag tecnici specifici in Italiano"],
  "tags_en": ["6-10 technical tags in English"],
  "technical_specs": ["strumenti/software/standard/algoritmi citati (se presenti)"]
}}

Regole:
- tags: stringhe brevi, NO frasi lunghe
- niente duplicati
- se non trovi technical_specs, usa [].

TESTO:
{chunk}
"""

            try:
                data = self._ollama_json_chat(map_prompt)

                summary = (data.get("summary") or "").strip()
                tags_it = self._normalize_tags_field(data.get("tags_it", []))
                tags_en = self._normalize_tags_field(data.get("tags_en", []))
                tech = self._normalize_tags_field(data.get("technical_specs", []))

                if summary:
                    partial_summaries.append(summary)

                partials.append({
                    "index": i,
                    "chunk_chars": len(chunk),
                    "summary": summary,
                    "tags_it": tags_it,
                    "tags_en": tags_en,
                    "technical_specs": tech,
                    "error": None,
                    "model": self.model,
                    "timestamp": datetime.datetime.now().isoformat(),
                })

                print(f"    - Chunk {i + 1}/{len(chunks)} analizzato.")
            except Exception as e:
                print(f"    [!] Errore chunk {i}: {e}")
                partials.append({
                    "index": i,
                    "chunk_chars": len(chunk),
                    "summary": "",
                    "tags_it": [],
                    "tags_en": [],
                    "technical_specs": [],
                    "error": str(e),
                    "model": self.model,
                    "timestamp": datetime.datetime.now().isoformat(),
                })

        if self.dump_chunks:
            self._save_chunks_json(partials)


        

        # 2b. Aggregazione locale Tag/Specs (union + ranking) dai chunk
        all_tags_it = [p.get("tags_it", []) for p in partials]
        all_tags_en = [p.get("tags_en", []) for p in partials]
        all_specs = [p.get("technical_specs", []) for p in partials]

        # Increase top_n to gather more candidates before deduplication
        ranked_tags_it = self._rank_union(all_tags_it, top_n=30)
        ranked_tags_en = self._rank_union(all_tags_en, top_n=30)
        ranked_specs = self._rank_union(all_specs, top_n=20)

        print("[*] Ottimizzazione e deduplicazione dei tag...")
        ranked_tags_it = self._deduplicate_tags_llm(ranked_tags_it, "it")[:10]
        ranked_tags_en = self._deduplicate_tags_llm(ranked_tags_en, "en")[:10]

        # 3. REDUCE (Sintesi Finale) con riduzione gerarchica per evitare superamento contesto
        print("[*] Generazione Abstract Finale e Tag...")

        def hierarchical_reduce(summaries):
            if not summaries:
                return {}

            group_size = 6
            groups = [summaries[i : i + group_size] for i in range(0, len(summaries), group_size)]

            condensed = []
            for gi, grp in enumerate(groups):
                grp_text = "\n---\n".join(grp)
                grp_prompt = f"""
                Riassumi questi riassunti intermedi (Gruppo {gi + 1}) in un paragrafo coeso in Inglese.
                OUTPUT JSON:
                {{"abstract": "..."}}

                INPUT:
                {grp_text}
                """
                data = self._ollama_json_chat(grp_prompt)
                if data.get("abstract"):
                    condensed.append(data["abstract"])
                else:
                    print(f"    [!] Warning: Intermediate reduction failed for group {gi+1}")
                    if self.debug_reduce:
                        # Try to capture raw output for debugging
                        try:
                            raw = ollama.chat(
                                model=self.model,
                                messages=[{"role": "user", "content": grp_prompt}],
                                options={"temperature": 0},
                            )
                            raw_text = raw.get("message", {}).get("content", "")
                            with open(f"{self.filename}.reduce_group_{gi}.raw.txt", "w", encoding="utf-8") as rf:
                                rf.write(raw_text)
                        except Exception:
                            pass

            final_input = "\n---\n".join(condensed) if condensed else "\n---\n".join(summaries)

            
            reduce_prompt = f"""
You are a technical analyst.
Below are partial summaries extracted from a technical document.

Task:
- Produce ONLY JSON with key "abstract" (English).
- The abstract must be cohesive and cover: scope, methodology (technical details), results, conclusions.

OUTPUT JSON:
{{"abstract":"..."}}

INPUT:
{final_input}
"""

            # First: ask for JSON directly

            direct = self._ollama_json_chat(reduce_prompt, options={"temperature": 0})
            if direct:
                return direct

            # Second: get raw text and try to extract JSON
            try:
                raw = ollama.chat(
                    model=self.model,
                    messages=[{"role": "user", "content": reduce_prompt}],
                    options={"temperature": 0},
                )
                raw_text = raw.get("message", {}).get("content", "")
            except Exception:
                raw_text = ""

            if self.debug_reduce and raw_text:
                with open(f"{self.filename}.reduce_raw.txt", "w", encoding="utf-8") as f:
                    f.write(raw_text)

            extracted = self._extract_json_object(raw_text)
            if extracted:
                return extracted

            return {}

        def fallback_reduce(summaries, fields_to_fetch=None):
            """Fallback: ask for fields one by one and merge."""
            base_text = "\n---\n".join(summaries)

            prompts = {
                "abstract": f"""
                Write ONLY JSON with key 'abstract' (English ~300-400 words) summarizing scope, methodology (technical details), results, conclusions.
                INPUT:
                {base_text}
                """,
                "tags_it": f"""
                Return ONLY JSON with key 'tags_it' as an array of 10 specific technical tags in Italian.
                INPUT:
                {base_text}
                """,
                "tags_en": f"""
                Return ONLY JSON with key 'tags_en' as an array of 10 specific technical tags in English.
                INPUT:
                {base_text}
                """,
                "technical_specs": f"""
                Return ONLY JSON with key 'technical_specs' as an array listing tools/software/algorithms cited.
                INPUT:
                {base_text}
                """
            }

            result = {}
            for name, prompt in prompts.items():
                if fields_to_fetch and name not in fields_to_fetch:
                    continue

                data = self._ollama_json_chat(prompt, options={"temperature": 0})
                if not data:
                    try:
                        raw = ollama.chat(model=self.model, messages=[{"role": "user", "content": prompt}], options={"temperature": 0})
                        raw_text = raw.get("message", {}).get("content", "")
                        data = self._extract_json_object(raw_text)
                    except Exception:
                        data = {}
                if name in data:
                    result[name] = data[name]

            return result

        try:
            data = hierarchical_reduce(partial_summaries)
            
            # Check for missing fields and try fallback
            required_fields = ["abstract"]
            missing_fields = [f for f in required_fields if f not in data or not data[f]]
            
            if missing_fields:
                print(f"[*] Missing fields after reduce: {missing_fields}. Attempting fallback...")
                fallback_data = fallback_reduce(partial_summaries, missing_fields)
                data.update(fallback_data)

            self.metadata["analysis"] = {
                "abstract": data.get("abstract", "N/A"),
                "tags": {"it": ranked_tags_it, "en": ranked_tags_en},
                "technical_specs": ranked_specs[:10],
                "method": "map_reduce",
                "chunks_processed": len(chunks),
                "word_count": len(self.content.split()),
                "model_used": self.model,
                "processed_date": datetime.datetime.now().isoformat(),
                "partials_dump": self.dump_chunks,
                "reduce_debug": self.debug_reduce,
                "partial_chunks": partials,
            }

            # Se l'output è vuoto, registra un avviso
            if self.metadata["analysis"]["abstract"] == "N/A" and not (
                self.metadata["analysis"]["tags"]["it"] or self.metadata["analysis"]["tags"]["en"]
            ):
                print("[!] Output JSON vuoto dal Reduce. Considera pull/aggiornare il modello o riprovare.")

        except Exception as e:
            print(f"[!] Errore nel Reduce finale: {e}")
            self.metadata["analysis"] = {
                "partial_summaries": partial_summaries,
                "error": str(e),
            }

    def _save_chunks_json(self, partials):
        output_filename = f"{self.filename}.chunks.json"
        with open(output_filename, "w", encoding="utf-8") as f:
            json.dump({
                "file": self.filename,
                "model": self.model,
                "chunks": partials,
            }, f, indent=4, ensure_ascii=False)
        print(f"[*] Dump dei chunk salvato in {output_filename}")

    def _smart_chunking(self, text, chunk_size=6000, overlap=500):
        """Divide il testo in chunk con sovrapposizione."""
        if len(text) <= chunk_size:
            return [text]

        chunks = []
        start = 0
        text_len = len(text)

        while start < text_len:
            end = start + chunk_size

            if end >= text_len:
                chunks.append(text[start:])
                break

            snippet = text[end - 100 : end]
            last_space = snippet.rfind(" ")

            if last_space != -1:
                end = (end - 100) + last_space

            chunks.append(text[start:end])

            start = end - overlap

        return chunks

    def _save_json(self):
        output_filename = f"{self.filename}.json"
        with open(output_filename, "w", encoding="utf-8") as f:
            json.dump(self.metadata, f, indent=4, ensure_ascii=False)


if __name__ == "__main__":
    parser = argparse.ArgumentParser(
        prog="doctagger.py",
        description="DocTagger – AI-powered document analysis using local Ollama models.",
        epilog="""
Examples:
  # Basic usage with default model (gemma3)
  python doctagger.py document.pdf
  
  # Use a different model
  python doctagger.py document.pdf --model mistral
  
  # Save per-chunk summaries
  python doctagger.py document.pdf --dump-chunks
  
  # Enable reduce debugging (for troubleshooting empty output)
  python doctagger.py document.pdf --debug-reduce
  
  # Combine multiple flags
  python doctagger.py document.pdf --model gemma3 --dump-chunks --debug-reduce

Supported formats: PDF, DOCX, PPTX, TXT

Output files:
  - <filename>.json – Main analysis with abstract, tags (IT/EN), technical specs, and partial chunks.
  - <filename>.chunks.json – Per-chunk summaries (with --dump-chunks).
  - <filename>.reduce_raw.txt – Raw reduce output (with --debug-reduce if JSON parsing failed).
  - <filename>.reduce_group_<N>.raw.txt – Intermediate group outputs (with --debug-reduce).

For more information, see README.md or run:
  python doctagger.py --help
        """,
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    parser.add_argument(
        "file",
        help="Path to the document to analyze (PDF, DOCX, PPTX, or TXT)."
    )
    parser.add_argument(
        "--model",
        default="gemma3",
        help="Ollama model to use for analysis. Default: gemma3. Examples: mistral, llama2, llama3.2."
    )
    parser.add_argument(
        "--dump-chunks",
        action="store_true",
        help="Save per-chunk analysis (index, char count, summary, error, timestamp) to <filename>.chunks.json."
    )
    parser.add_argument(
        "--debug-reduce",
        action="store_true",
        help="Dump raw reduce outputs and enable JSON extraction fallback. Writes <filename>.reduce_raw.txt and optionally <filename>.reduce_group_*.raw.txt."
    )
    args = parser.parse_args()

    if os.path.exists(args.file):
        analyzer = DocumentAnalyzer(
            args.file,
            model=args.model,
            dump_chunks=args.dump_chunks,
            debug_reduce=args.debug_reduce,
        )
        analyzer.process()
    else:
        print("File not found.")